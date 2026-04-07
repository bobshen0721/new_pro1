import base64
import io
import math
import mimetypes
import os
import tempfile
from pathlib import Path

import gradio as gr
import anthropic
import google.generativeai as genai
from openai import OpenAI
from pydub import AudioSegment

# ---------------------------------------------------------------------------
# File processing utilities
# ---------------------------------------------------------------------------

SUPPORTED_EXTENSIONS = {
    "image": [".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".svg"],
    "pdf": [".pdf"],
    "document": [".txt", ".md", ".csv", ".json", ".xml", ".html", ".htm",
                 ".py", ".js", ".ts", ".java", ".c", ".cpp", ".r", ".sql"],
    "excel": [".xlsx", ".xls"],
    "ppt": [".pptx", ".ppt"],
    "audio": [".mp3", ".wav", ".m4a", ".ogg", ".flac", ".aac", ".wma", ".webm"],
}

AUDIO_CHUNK_MS = 30_000  # 30 seconds per chunk, same as Whisper


def _file_category(filepath: str) -> str:
    ext = Path(filepath).suffix.lower()
    for cat, exts in SUPPORTED_EXTENSIONS.items():
        if ext in exts:
            return cat
    return "unknown"


def _read_bytes(filepath: str) -> bytes:
    with open(filepath, "rb") as f:
        return f.read()


def _read_text(filepath: str) -> str:
    with open(filepath, "r", errors="replace") as f:
        return f.read()


def _extract_excel(filepath: str) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            rows.append("\t".join(str(c) if c is not None else "" for c in row))
        parts.append(f"=== Sheet: {sheet_name} ===\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(parts)


def _extract_pptx(filepath: str) -> str:
    from pptx import Presentation
    prs = Presentation(filepath)
    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        slides_text.append(f"=== Slide {i} ===\n" + "\n".join(texts))
    return "\n\n".join(slides_text)


def _extract_pdf_text(filepath: str) -> str:
    import pymupdf
    doc = pymupdf.open(filepath)
    pages = []
    for i, page in enumerate(doc, 1):
        pages.append(f"=== Page {i} ===\n{page.get_text()}")
    doc.close()
    return "\n\n".join(pages)


# ---------------------------------------------------------------------------
# Audio chunking & transcription (unlimited length, 30s chunks)
# ---------------------------------------------------------------------------

def _split_audio_chunks(filepath: str, chunk_ms: int = AUDIO_CHUNK_MS) -> list[AudioSegment]:
    """Split an audio file into fixed-length chunks."""
    audio = AudioSegment.from_file(filepath)
    chunks = []
    for start in range(0, len(audio), chunk_ms):
        chunks.append(audio[start:start + chunk_ms])
    return chunks


def _chunk_to_bytes(chunk: AudioSegment, fmt: str = "mp3") -> bytes:
    """Export an AudioSegment chunk to bytes."""
    buf = io.BytesIO()
    chunk.export(buf, format=fmt)
    return buf.getvalue()


def transcribe_audio_gemini(filepath: str, api_key: str, model: str,
                            language_hint: str = "") -> str:
    """
    Transcribe an audio file of ANY length using Gemini.
    Splits into 30-second chunks and processes each one sequentially,
    just like Whisper does internally.
    """
    genai.configure(api_key=api_key)
    gen_model = genai.GenerativeModel(model_name=model)

    chunks = _split_audio_chunks(filepath)
    total = len(chunks)
    audio_duration = AudioSegment.from_file(filepath).duration_seconds

    transcription_parts = []

    lang_instruction = f" The audio language is {language_hint}." if language_hint else ""
    for i, chunk in enumerate(chunks):
        audio_bytes = _chunk_to_bytes(chunk, fmt="mp3")

        prompt = (
            f"Transcribe the following audio accurately and completely. "
            f"This is chunk {i + 1}/{total} of a longer audio "
            f"(total duration: {audio_duration:.1f}s).{lang_instruction} "
            f"Output ONLY the transcribed text, no timestamps, no commentary."
        )

        response = gen_model.generate_content([
            prompt,
            {"mime_type": "audio/mp3", "data": audio_bytes},
        ])

        text = response.text.strip() if response.text else ""
        if text:
            transcription_parts.append(text)

    return "\n".join(transcription_parts)


def transcribe_audio_openai(filepath: str, api_key: str, model: str,
                            language_hint: str = "") -> str:
    """
    Transcribe an audio file of ANY length using OpenAI Whisper API.
    Splits into 30-second chunks for files that exceed API limits.
    """
    client = OpenAI(api_key=api_key)
    chunks = _split_audio_chunks(filepath)
    transcription_parts = []

    for i, chunk in enumerate(chunks):
        audio_bytes = _chunk_to_bytes(chunk, fmt="mp3")
        audio_file = io.BytesIO(audio_bytes)
        audio_file.name = f"chunk_{i}.mp3"

        kwargs = {"model": model, "file": audio_file}
        if language_hint:
            kwargs["language"] = language_hint

        result = client.audio.transcriptions.create(**kwargs)
        text = result.text.strip() if result.text else ""
        if text:
            transcription_parts.append(text)

    return "\n".join(transcription_parts)


def transcribe_audio_streaming(filepath: str, api_key: str, model: str,
                               provider: str, language_hint: str = ""):
    """
    Generator that yields progressive transcription results.
    Shows real-time progress as each chunk completes.
    """
    audio = AudioSegment.from_file(filepath)
    total_duration = audio.duration_seconds
    chunks = _split_audio_chunks(filepath)
    total = len(chunks)

    header = (
        f"🎙️ Audio: {Path(filepath).name}\n"
        f"⏱️ Duration: {total_duration:.1f}s | "
        f"📦 Chunks: {total} (30s each)\n"
        f"🔄 Model: {model}\n\n"
        f"---\n\n"
    )

    transcription_parts = []
    progress_text = header + "Transcribing...\n\n"
    yield progress_text

    if provider == "Gemini":
        genai.configure(api_key=api_key)
        gen_model = genai.GenerativeModel(model_name=model)

        lang_instruction = f" The audio language is {language_hint}." if language_hint else ""
        for i, chunk in enumerate(chunks):
            audio_bytes = _chunk_to_bytes(chunk, fmt="mp3")
            chunk_start = i * 30
            chunk_end = min((i + 1) * 30, total_duration)

            prompt = (
                f"Transcribe the following audio accurately and completely. "
                f"This is chunk {i + 1}/{total} of a longer audio "
                f"(total duration: {total_duration:.1f}s).{lang_instruction} "
                f"Output ONLY the transcribed text, no timestamps, no commentary."
            )

            response = gen_model.generate_content([
                prompt,
                {"mime_type": "audio/mp3", "data": audio_bytes},
            ])

            text = response.text.strip() if response.text else ""
            if text:
                transcription_parts.append(text)

            progress = (i + 1) / total * 100
            progress_text = (
                header
                + f"Progress: {i + 1}/{total} chunks ({progress:.0f}%) "
                + f"[{chunk_start:.0f}s - {chunk_end:.0f}s]\n\n"
                + "---\n\n"
                + "\n".join(transcription_parts)
            )
            yield progress_text

    elif provider == "OpenAI":
        client = OpenAI(api_key=api_key)
        for i, chunk in enumerate(chunks):
            audio_bytes = _chunk_to_bytes(chunk, fmt="mp3")
            audio_file = io.BytesIO(audio_bytes)
            audio_file.name = f"chunk_{i}.mp3"
            chunk_start = i * 30
            chunk_end = min((i + 1) * 30, total_duration)

            kwargs = {"model": model, "file": audio_file}
            if language_hint:
                kwargs["language"] = language_hint
            result = client.audio.transcriptions.create(**kwargs)
            text = result.text.strip() if result.text else ""
            if text:
                transcription_parts.append(text)

            progress = (i + 1) / total * 100
            progress_text = (
                header
                + f"Progress: {i + 1}/{total} chunks ({progress:.0f}%) "
                + f"[{chunk_start:.0f}s - {chunk_end:.0f}s]\n\n"
                + "---\n\n"
                + "\n".join(transcription_parts)
            )
            yield progress_text

    # Final output
    final = (
        header
        + f"✅ Completed! {total} chunks transcribed.\n\n"
        + "---\n\n"
        + "\n".join(transcription_parts)
    )
    yield final


# ---------------------------------------------------------------------------
# Provider: OpenAI
# ---------------------------------------------------------------------------

def _build_openai_file_content(filepath: str):
    """Return a list of content parts for OpenAI multi-modal messages."""
    category = _file_category(filepath)
    mime = mimetypes.guess_type(filepath)[0] or "application/octet-stream"

    if category == "image":
        b64 = base64.b64encode(_read_bytes(filepath)).decode()
        return [{"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64}"}}]
    elif category == "pdf":
        text = _extract_pdf_text(filepath)
        return [{"type": "text", "text": f"[PDF content from {Path(filepath).name}]\n{text}"}]
    elif category == "excel":
        text = _extract_excel(filepath)
        return [{"type": "text", "text": f"[Excel content from {Path(filepath).name}]\n{text}"}]
    elif category == "ppt":
        text = _extract_pptx(filepath)
        return [{"type": "text", "text": f"[PPT content from {Path(filepath).name}]\n{text}"}]
    elif category == "document":
        text = _read_text(filepath)
        return [{"type": "text", "text": f"[File: {Path(filepath).name}]\n{text}"}]
    else:
        return [{"type": "text", "text": f"[Unsupported file: {Path(filepath).name}]"}]


def chat_openai(message: str, files: list[str], history: list, model: str, api_key: str, system_prompt: str):
    client = OpenAI(api_key=api_key)
    messages = [{"role": "system", "content": system_prompt}]

    for msg in history:
        messages.append(msg)

    # Build current user message
    content_parts = [{"type": "text", "text": message}] if message else []
    for fp in files:
        content_parts.extend(_build_openai_file_content(fp))

    messages.append({"role": "user", "content": content_parts if len(content_parts) > 1 else message or "Please analyze the uploaded files."})

    response = client.chat.completions.create(model=model, messages=messages, stream=True)
    partial = ""
    for chunk in response:
        delta = chunk.choices[0].delta.content
        if delta:
            partial += delta
            yield partial


# ---------------------------------------------------------------------------
# Provider: Anthropic (Claude)
# ---------------------------------------------------------------------------

def _build_claude_file_content(filepath: str):
    category = _file_category(filepath)
    mime = mimetypes.guess_type(filepath)[0] or "application/octet-stream"

    if category == "image":
        b64 = base64.b64encode(_read_bytes(filepath)).decode()
        return [{"type": "image", "source": {"type": "base64", "media_type": mime, "data": b64}}]
    elif category == "pdf":
        b64 = base64.b64encode(_read_bytes(filepath)).decode()
        return [{"type": "document", "source": {"type": "base64", "media_type": "application/pdf", "data": b64}}]
    elif category == "excel":
        text = _extract_excel(filepath)
        return [{"type": "text", "text": f"[Excel content from {Path(filepath).name}]\n{text}"}]
    elif category == "ppt":
        text = _extract_pptx(filepath)
        return [{"type": "text", "text": f"[PPT content from {Path(filepath).name}]\n{text}"}]
    elif category == "document":
        text = _read_text(filepath)
        return [{"type": "text", "text": f"[File: {Path(filepath).name}]\n{text}"}]
    else:
        return [{"type": "text", "text": f"[Unsupported file: {Path(filepath).name}]"}]


def chat_claude(message: str, files: list[str], history: list, model: str, api_key: str, system_prompt: str):
    client = anthropic.Anthropic(api_key=api_key)

    messages = []
    for msg in history:
        messages.append(msg)

    content_parts = []
    for fp in files:
        content_parts.extend(_build_claude_file_content(fp))
    content_parts.append({"type": "text", "text": message or "Please analyze the uploaded files."})

    messages.append({"role": "user", "content": content_parts})

    with client.messages.stream(model=model, max_tokens=8192, system=system_prompt, messages=messages) as stream:
        partial = ""
        for text in stream.text_stream:
            partial += text
            yield partial


# ---------------------------------------------------------------------------
# Provider: Google Gemini
# ---------------------------------------------------------------------------

def _build_gemini_file_content(filepath: str):
    category = _file_category(filepath)
    mime = mimetypes.guess_type(filepath)[0] or "application/octet-stream"

    if category == "image":
        return [{"mime_type": mime, "data": _read_bytes(filepath)}]
    elif category == "audio":
        return [{"mime_type": mime, "data": _read_bytes(filepath)}]
    elif category == "pdf":
        text = _extract_pdf_text(filepath)
        return [f"[PDF content from {Path(filepath).name}]\n{text}"]
    elif category == "excel":
        text = _extract_excel(filepath)
        return [f"[Excel content from {Path(filepath).name}]\n{text}"]
    elif category == "ppt":
        text = _extract_pptx(filepath)
        return [f"[PPT content from {Path(filepath).name}]\n{text}"]
    elif category == "document":
        text = _read_text(filepath)
        return [f"[File: {Path(filepath).name}]\n{text}"]
    else:
        return [f"[Unsupported file: {Path(filepath).name}]"]


def chat_gemini(message: str, files: list[str], history: list, model: str, api_key: str, system_prompt: str):
    genai.configure(api_key=api_key)
    gen_model = genai.GenerativeModel(model_name=model, system_instruction=system_prompt)

    # Build Gemini history
    gemini_history = []
    for msg in history:
        role = "user" if msg["role"] == "user" else "model"
        gemini_history.append({"role": role, "parts": [msg["content"]] if isinstance(msg["content"], str) else [msg["content"]]})

    chat_session = gen_model.start_chat(history=gemini_history)

    parts = []
    for fp in files:
        parts.extend(_build_gemini_file_content(fp))
    parts.append(message or "Please analyze the uploaded files.")

    response = chat_session.send_message(parts, stream=True)
    partial = ""
    for chunk in response:
        if chunk.text:
            partial += chunk.text
            yield partial


# ---------------------------------------------------------------------------
# Provider: xAI Grok (OpenAI-compatible API)
# ---------------------------------------------------------------------------

def chat_grok(message: str, files: list[str], history: list, model: str, api_key: str, system_prompt: str):
    client = OpenAI(api_key=api_key, base_url="https://api.x.ai/v1")
    messages = [{"role": "system", "content": system_prompt}]

    for msg in history:
        messages.append(msg)

    content_parts = [{"type": "text", "text": message}] if message else []
    for fp in files:
        content_parts.extend(_build_openai_file_content(fp))

    messages.append({"role": "user", "content": content_parts if len(content_parts) > 1 else message or "Please analyze the uploaded files."})

    response = client.chat.completions.create(model=model, messages=messages, stream=True)
    partial = ""
    for chunk in response:
        delta = chunk.choices[0].delta.content
        if delta:
            partial += delta
            yield partial


# ---------------------------------------------------------------------------
# Available models per provider
# ---------------------------------------------------------------------------

PROVIDER_MODELS = {
    "OpenAI": ["gpt-5.3", "gpt-5.3-codex", "gpt-4o", "gpt-4o-mini", "gpt-4.1", "gpt-4.1-mini", "gpt-4.1-nano", "o3-mini"],
    "Claude": ["claude-sonnet-4-20250514", "claude-haiku-4-5-20251001", "claude-opus-4-20250514"],
    "Gemini": ["gemini-3.1-pro", "gemini-3.1-flash", "gemini-3.1-flash-lite", "gemini-3-pro", "gemini-3-flash", "gemini-2.5-pro", "gemini-2.5-flash", "gemini-2.0-flash"],
    "Grok": ["grok-4.1-fast", "grok-4.20", "grok-code-fast-1", "grok-beta"],
}

PROVIDER_FN = {
    "OpenAI": chat_openai,
    "Claude": chat_claude,
    "Gemini": chat_gemini,
    "Grok": chat_grok,
}


# ---------------------------------------------------------------------------
# Gradio UI
# ---------------------------------------------------------------------------

def update_model_choices(provider):
    models = PROVIDER_MODELS.get(provider, [])
    return gr.Dropdown(choices=models, value=models[0] if models else "")


def respond(message, chat_history, uploaded_files, provider, model, api_key, system_prompt):
    if not api_key:
        yield chat_history + [{"role": "assistant", "content": f"⚠️ Please enter your {provider} API Key in the sidebar."}]
        return

    # Collect file paths
    file_paths = []
    if uploaded_files:
        for f in uploaded_files:
            file_paths.append(f.name if hasattr(f, "name") else str(f))

    # Build display message with file indicators
    display_parts = []
    if message:
        display_parts.append(message)
    for fp in file_paths:
        display_parts.append(f"📎 {Path(fp).name}")
    display_message = "\n".join(display_parts) or "📎 (files uploaded)"

    chat_history = chat_history + [{"role": "user", "content": display_message}]
    yield chat_history

    # Build API history (without the current message)
    api_history = []
    for msg in chat_history[:-1]:
        api_history.append({"role": msg["role"], "content": msg["content"]})

    chat_fn = PROVIDER_FN[provider]

    try:
        for partial_response in chat_fn(message, file_paths, api_history, model, api_key, system_prompt):
            yield chat_history + [{"role": "assistant", "content": partial_response}]
    except Exception as e:
        error_msg = f"❌ Error: {str(e)}"
        yield chat_history + [{"role": "assistant", "content": error_msg}]


def clear_chat():
    return [], None


TRANSCRIPTION_MODELS = {
    "Gemini": ["gemini-2.5-flash", "gemini-2.5-pro", "gemini-2.0-flash",
               "gemini-3.1-flash", "gemini-3.1-pro"],
    "OpenAI": ["whisper-1", "gpt-4o-transcribe", "gpt-4o-mini-transcribe"],
}


def update_transcription_models(provider):
    models = TRANSCRIPTION_MODELS.get(provider, [])
    return gr.Dropdown(choices=models, value=models[0] if models else "")


def run_transcription(audio_file, provider, model, api_key, language_hint):
    if not audio_file:
        yield "Please upload an audio file."
        return
    if not api_key:
        yield f"Please enter your {provider} API Key."
        return

    filepath = audio_file.name if hasattr(audio_file, "name") else str(audio_file)

    try:
        for progress in transcribe_audio_streaming(
            filepath, api_key, model, provider, language_hint
        ):
            yield progress
    except Exception as e:
        yield f"❌ Error: {str(e)}"


# Build the interface
with gr.Blocks(theme=gr.themes.Soft(), title="Multi-AI Chatbot") as demo:
    gr.Markdown("# 🤖 Multi-AI Chatbot\nChat with **OpenAI**, **Claude**, **Gemini**, or **Grok** — supports file uploads and unlimited audio transcription.")

    with gr.Tabs():
        # ============== Tab 1: Chat ==============
        with gr.Tab("💬 Chat"):
            with gr.Row():
                # Sidebar
                with gr.Column(scale=1, min_width=280):
                    provider = gr.Dropdown(
                        choices=["OpenAI", "Claude", "Gemini", "Grok"],
                        value="OpenAI",
                        label="AI Provider",
                    )
                    model = gr.Dropdown(
                        choices=PROVIDER_MODELS["OpenAI"],
                        value=PROVIDER_MODELS["OpenAI"][0],
                        label="Model",
                    )
                    api_key = gr.Textbox(
                        label="API Key",
                        type="password",
                        placeholder="Enter your API key...",
                    )
                    system_prompt = gr.Textbox(
                        label="System Prompt",
                        value="You are a helpful assistant. Answer in the same language the user uses.",
                        lines=3,
                    )
                    uploaded_files = gr.File(
                        label="Upload Files",
                        file_count="multiple",
                        file_types=[
                            ".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp",
                            ".pdf", ".txt", ".md", ".csv", ".json", ".xml",
                            ".html", ".py", ".js", ".ts", ".java", ".c", ".cpp",
                            ".r", ".sql", ".xlsx", ".xls", ".pptx", ".ppt",
                            ".mp3", ".wav", ".m4a", ".ogg", ".flac", ".aac",
                        ],
                    )
                    gr.Markdown(
                        "**Supported files:** Images, PDF, TXT, CSV, JSON, "
                        "Excel (.xlsx), PowerPoint (.pptx), audio files, code files, and more."
                    )

                # Chat area
                with gr.Column(scale=3):
                    chatbot = gr.Chatbot(
                        height=600,
                        type="messages",
                        show_copy_button=True,
                        placeholder="Select a provider, enter your API key, and start chatting!",
                    )
                    with gr.Row():
                        msg = gr.Textbox(
                            placeholder="Type your message here...",
                            show_label=False,
                            scale=6,
                            container=False,
                        )
                        send_btn = gr.Button("Send", variant="primary", scale=1)
                        clear_btn = gr.Button("Clear", scale=1)

        # ============== Tab 2: Audio Transcription ==============
        with gr.Tab("🎙️ Audio Transcription"):
            gr.Markdown(
                "### Unlimited Audio Transcription\n"
                "Upload audio of **any length**. The file is split into 30-second chunks "
                "and transcribed sequentially — just like Whisper does internally.\n\n"
                "Supports: MP3, WAV, M4A, OGG, FLAC, AAC, WMA, WebM"
            )
            with gr.Row():
                with gr.Column(scale=1, min_width=280):
                    stt_provider = gr.Dropdown(
                        choices=["Gemini", "OpenAI"],
                        value="Gemini",
                        label="Transcription Provider",
                    )
                    stt_model = gr.Dropdown(
                        choices=TRANSCRIPTION_MODELS["Gemini"],
                        value=TRANSCRIPTION_MODELS["Gemini"][0],
                        label="Model",
                    )
                    stt_api_key = gr.Textbox(
                        label="API Key",
                        type="password",
                        placeholder="Enter your API key...",
                    )
                    stt_language = gr.Textbox(
                        label="Language Hint (optional)",
                        placeholder="e.g. zh, en, ja, ko...",
                    )
                    stt_audio = gr.File(
                        label="Upload Audio File",
                        file_count="single",
                        file_types=[
                            ".mp3", ".wav", ".m4a", ".ogg", ".flac",
                            ".aac", ".wma", ".webm",
                        ],
                    )
                    stt_btn = gr.Button("🎙️ Start Transcription", variant="primary")

                with gr.Column(scale=3):
                    stt_output = gr.Textbox(
                        label="Transcription Result",
                        lines=25,
                        show_copy_button=True,
                        interactive=False,
                    )

    # ============== Events ==============
    # Chat tab
    provider.change(fn=update_model_choices, inputs=provider, outputs=model)

    send_btn.click(
        fn=respond,
        inputs=[msg, chatbot, uploaded_files, provider, model, api_key, system_prompt],
        outputs=chatbot,
    ).then(lambda: ("", None), outputs=[msg, uploaded_files])

    msg.submit(
        fn=respond,
        inputs=[msg, chatbot, uploaded_files, provider, model, api_key, system_prompt],
        outputs=chatbot,
    ).then(lambda: ("", None), outputs=[msg, uploaded_files])

    clear_btn.click(fn=clear_chat, outputs=[chatbot, uploaded_files])

    # Audio transcription tab
    stt_provider.change(fn=update_transcription_models, inputs=stt_provider, outputs=stt_model)

    stt_btn.click(
        fn=run_transcription,
        inputs=[stt_audio, stt_provider, stt_model, stt_api_key, stt_language],
        outputs=stt_output,
    )

if __name__ == "__main__":
    demo.launch()
